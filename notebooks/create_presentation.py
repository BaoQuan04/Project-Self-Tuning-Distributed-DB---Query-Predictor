"""
create_presentation.py
Tạo file PowerPoint: results/Query_Predictor_Presentation.pptx
Chạy: python notebooks/create_presentation.py
"""
import os, sys, json
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as pns
from lxml import etree

# ── Colors ────────────────────────────────────────────────────────────────────
C_DARK   = RGBColor(0x1A, 0x29, 0x4A)   # navy
C_BLUE   = RGBColor(0x21, 0x96, 0xF3)   # accent blue
C_GREEN  = RGBColor(0x2E, 0xCC, 0x71)   # green
C_RED    = RGBColor(0xE7, 0x4C, 0x3C)   # red
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY   = RGBColor(0x55, 0x55, 0x55)
C_LIGHT  = RGBColor(0xF0, 0xF4, 0xF8)

CHART_DIR   = os.path.join('results', 'charts')
OUTPUT_PATH = os.path.join('results', 'Query_Predictor_Presentation.pptx')
RESULTS_PATH = os.path.join('results', 'reports', 'benchmark_summary.json')

W, H = Inches(13.33), Inches(7.5)   # 16:9 widescreen


# ── Helpers ───────────────────────────────────────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)


def add_rect(slide, left, top, width, height, fill_rgb, alpha=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=C_WHITE, align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size    = Pt(font_size)
    run.font.bold    = bold
    run.font.italic  = italic
    run.font.color.rgb = color
    return txb


def add_image(slide, path, left, top, width, height=None):
    if not os.path.exists(path):
        return
    if height:
        slide.shapes.add_picture(path, left, top, width, height)
    else:
        slide.shapes.add_picture(path, left, top, width)


def header_bar(slide, title, subtitle=None):
    add_rect(slide, 0, 0, W, Inches(1.15), C_DARK)
    add_text(slide, title, Inches(0.4), Inches(0.1), Inches(12), Inches(0.65),
             font_size=28, bold=True, color=C_WHITE)
    if subtitle:
        add_text(slide, subtitle, Inches(0.4), Inches(0.72), Inches(12), Inches(0.4),
                 font_size=15, color=C_BLUE)


def bullet_box(slide, items, left, top, width, height,
               font_size=16, color=C_DARK, bullet='•'):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f'{bullet}  {item}'
        run.font.size = Pt(font_size)
        run.font.color.rgb = color


def stat_box(slide, value, label, left, top, bg_color):
    add_rect(slide, left, top, Inches(2.8), Inches(1.5), bg_color)
    add_text(slide, value, left+Inches(0.1), top+Inches(0.1),
             Inches(2.6), Inches(0.8), font_size=34, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, label, left+Inches(0.1), top+Inches(0.85),
             Inches(2.6), Inches(0.55), font_size=13, color=C_WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDES
# ══════════════════════════════════════════════════════════════════════════════

def slide_title(prs):
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, H, C_DARK)
    # accent stripe
    add_rect(s, 0, Inches(3.6), W, Inches(0.06), C_BLUE)

    add_text(s, 'QUERY PREDICTOR', Inches(0.6), Inches(0.8), Inches(12), Inches(1.1),
             font_size=52, bold=True, color=C_WHITE)
    add_text(s, 'Self-Tuning Distributed Database with ML Pre-fetching',
             Inches(0.6), Inches(1.85), Inches(11), Inches(0.7),
             font_size=24, color=C_BLUE)
    add_text(s, 'Topic 142  |  Distributed Database Systems',
             Inches(0.6), Inches(2.55), Inches(11), Inches(0.5),
             font_size=18, italic=True, color=RGBColor(0xAA, 0xC4, 0xE0))

    # Tech badges
    badges = [('PostgreSQL ×2', C_BLUE), ('Redis Cache', C_GREEN),
              ('Markov Chain', RGBColor(0xF3,0x9C,0x12)), ('LSTM (PyTorch)', C_RED)]
    for i, (label, color) in enumerate(badges):
        bx = Inches(0.6) + i * Inches(3.0)
        add_rect(s, bx, Inches(4.1), Inches(2.7), Inches(0.55), color)
        add_text(s, label, bx+Inches(0.1), Inches(4.12), Inches(2.5), Inches(0.5),
                 font_size=15, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    add_text(s, '2026  |  Distributed Systems Project',
             Inches(0.6), Inches(6.8), Inches(12), Inches(0.4),
             font_size=13, color=RGBColor(0x88, 0x99, 0xAA), align=PP_ALIGN.LEFT)


def slide_problem(prs):
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, H, C_LIGHT)
    header_bar(s, 'Problem Statement', 'Why does distributed DB querying suffer?')

    problems = [
        'Cross-site queries (Site A → Site B) incur 150ms+ network latency',
        'No intelligence — every cache miss triggers a full remote DB call',
        'Repeated query patterns waste bandwidth fetching the same data repeatedly',
        'Result: slow response times and poor user experience at scale',
    ]
    bullet_box(s, problems, Inches(0.5), Inches(1.4), Inches(6.2), Inches(3.5),
               font_size=17, color=C_DARK)

    add_rect(s, Inches(7.0), Inches(1.3), Inches(5.8), Inches(4.2), C_DARK)
    add_text(s, 'Our Solution', Inches(7.1), Inches(1.4), Inches(5.6), Inches(0.5),
             font_size=18, bold=True, color=C_BLUE)
    solutions = [
        'Train ML model on query history',
        'Predict next query BEFORE it arrives',
        'Pre-fetch data into Redis cache',
        'Achieve cache hits instead of remote DB calls',
    ]
    bullet_box(s, solutions, Inches(7.1), Inches(1.9), Inches(5.5), Inches(2.8),
               font_size=16, color=C_WHITE)
    add_text(s, 'Expected: Cache Hit Rate 0% → 60-80%',
             Inches(7.1), Inches(4.8), Inches(5.5), Inches(0.5),
             font_size=14, bold=True, color=C_GREEN)


def slide_architecture(prs):
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, H, C_LIGHT)
    header_bar(s, 'System Architecture', 'Two-site distributed setup with ML prediction layer')

    components = [
        ('Site A  (Port 5435)', 'customers, orders\n→ Local queries (~5ms)', Inches(0.3), Inches(1.5), C_BLUE),
        ('Site B  (Port 5433)', 'products, inventory\n→ Remote queries (~150ms)', Inches(0.3), Inches(3.4), C_RED),
        ('Redis Cache', 'Shared in-memory cache\nTTL = 300s', Inches(4.7), Inches(2.4), C_GREEN),
        ('ML Predictor', 'Markov Chain (order=2)\n+  LSTM (2-layer)', Inches(9.0), Inches(1.5), RGBColor(0xF3,0x9C,0x12)),
        ('Pre-fetch Ctrl', 'Async background thread\nConfidence threshold ≥ 0.6', Inches(9.0), Inches(3.4), RGBColor(0x9B,0x59,0xB6)),
    ]
    for title, desc, x, y, color in components:
        add_rect(s, x, y, Inches(3.8), Inches(1.5), color)
        add_text(s, title, x+Inches(0.1), y+Inches(0.05), Inches(3.6), Inches(0.55),
                 font_size=15, bold=True, color=C_WHITE)
        add_text(s, desc, x+Inches(0.1), y+Inches(0.6), Inches(3.6), Inches(0.85),
                 font_size=12, color=C_WHITE)

    # Flow arrows as text
    add_text(s, '→  Cache hit / miss', Inches(4.2), Inches(1.9), Inches(0.9), Inches(0.8),
             font_size=10, color=C_GRAY)
    add_text(s, '↑  Pre-fetch data', Inches(4.2), Inches(3.7), Inches(0.9), Inches(0.8),
             font_size=10, color=C_GRAY)
    add_text(s, '← Predict', Inches(8.3), Inches(2.3), Inches(0.9), Inches(0.5),
             font_size=10, color=C_GRAY)

    add_text(s, 'Query Flow:  Client → Query Engine → Cache Check → (miss) → DB → Cache Store  →  Next: Pre-fetch predicted table',
             Inches(0.3), Inches(6.6), Inches(12.5), Inches(0.5),
             font_size=12, italic=True, color=C_GRAY)


def slide_dataset(prs):
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, H, C_LIGHT)
    header_bar(s, 'Dataset: Query_History_Logs', '500,000 synthetic query records — 3 months of history')

    stats = [
        ('500K', 'Total Query\nLogs', Inches(0.3)),
        ('1,000', 'Unique\nUsers', Inches(3.5)),
        ('6', 'Session\nPatterns', Inches(6.7)),
        ('4', 'Tables\n(2 per site)', Inches(9.9)),
    ]
    for val, label, x in stats:
        stat_box(s, val, label, x, Inches(1.5), C_DARK)

    add_text(s, 'Session Patterns (weighted):', Inches(0.5), Inches(3.4), Inches(5.5), Inches(0.4),
             font_size=15, bold=True, color=C_DARK)
    patterns = [
        'browse (30%):          SELECT products → inventory → orders',
        'purchase (25%):        SELECT products → inventory → INSERT orders → UPDATE inventory',
        'admin_report (15%):    SELECT orders → customers → products',
        'inventory_mgmt (12%):  SELECT inventory → UPDATE inventory',
        'customer_service (10%): SELECT customers → orders → UPDATE orders',
        'analytics (8%):        SELECT orders → customers → products',
    ]
    bullet_box(s, patterns, Inches(0.5), Inches(3.85), Inches(12.5), Inches(3.0),
               font_size=13, color=C_DARK, bullet='►')


def slide_ml_models(prs):
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, H, C_LIGHT)
    header_bar(s, 'ML Models: Next Query Prediction', 'Markov Chain (fast) + LSTM (accurate)')

    # Markov box
    add_rect(s, Inches(0.3), Inches(1.4), Inches(6.0), Inches(5.6), C_DARK)
    add_text(s, 'Markov Chain  (Order = 2)', Inches(0.4), Inches(1.5), Inches(5.8), Inches(0.55),
             font_size=18, bold=True, color=C_BLUE)
    markov_points = [
        'State = "QUERY_TYPE:table_name"',
        'Looks at last 2 queries as context',
        'Builds transition probability matrix',
        'Fast inference — O(1) lookup',
        'Best for: short sequential patterns',
        'Training: ~5 seconds',
    ]
    bullet_box(s, markov_points, Inches(0.4), Inches(2.1), Inches(5.7), Inches(3.5),
               font_size=14, color=C_WHITE)
    add_text(s, 'predict_top_k(["SELECT:products",\n  "SELECT:inventory"], k=3)',
             Inches(0.4), Inches(5.3), Inches(5.7), Inches(1.0),
             font_size=12, italic=True, color=RGBColor(0xAA,0xDD,0xFF))

    # LSTM box
    add_rect(s, Inches(6.9), Inches(1.4), Inches(6.0), Inches(5.6), RGBColor(0x1A,0x3A,0x2A))
    add_text(s, 'LSTM  (2-layer, hidden=64)', Inches(7.0), Inches(1.5), Inches(5.8), Inches(0.55),
             font_size=18, bold=True, color=C_GREEN)
    lstm_points = [
        'Input: last 5 query states (seq_len=5)',
        'Architecture: Embedding(16,32) → LSTM(64)×2 → Dense(16) → Softmax',
        'Captures long-range session dependencies',
        'Handles noise and irregular patterns',
        'Training: ~30 epochs, ~5-10 minutes',
        'Best for: complex multi-step workflows',
    ]
    bullet_box(s, lstm_points, Inches(7.0), Inches(2.1), Inches(5.8), Inches(3.5),
               font_size=14, color=C_WHITE)
    add_text(s, 'Confidence threshold: 0.6\n→ Only pre-fetch if model is 60%+ confident',
             Inches(7.0), Inches(5.5), Inches(5.8), Inches(1.0),
             font_size=12, italic=True, color=RGBColor(0x88,0xFF,0xBB))


def slide_benchmark_setup(prs):
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, H, C_LIGHT)
    header_bar(s, 'Benchmark Setup', '10,000 queries × 3 scenarios on identical workload')

    scenarios = [
        ('BASELINE', 'No cache, no prediction', 'Every query hits DB directly.\nSite B queries: 150ms latency.', C_RED, Inches(0.3)),
        ('CACHE ONLY', 'Redis cache, no ML', 'Cache miss → fetch from DB → store in cache.\nSubsequent identical queries: cache hit.', C_BLUE, Inches(4.6)),
        ('AI PRE-FETCH', 'Cache + ML prediction', 'After each query, ML predicts next.\nPre-fetches data into cache asynchronously.', C_GREEN, Inches(8.9)),
    ]
    for title, subtitle, desc, color, x in scenarios:
        add_rect(s, x, Inches(1.5), Inches(4.0), Inches(4.5), color)
        add_text(s, title, x+Inches(0.15), Inches(1.6), Inches(3.7), Inches(0.6),
                 font_size=22, bold=True, color=C_WHITE)
        add_text(s, subtitle, x+Inches(0.15), Inches(2.25), Inches(3.7), Inches(0.45),
                 font_size=13, italic=True, color=C_WHITE)
        add_text(s, desc, x+Inches(0.15), Inches(2.75), Inches(3.7), Inches(2.0),
                 font_size=14, color=C_WHITE)

    add_text(s, 'Same 10K test queries used for all 3 scenarios  |  Cache flushed between scenarios  |  LSTM predictor used for AI Pre-fetch',
             Inches(0.3), Inches(6.5), Inches(12.7), Inches(0.5),
             font_size=12, italic=True, color=C_GRAY, align=PP_ALIGN.CENTER)


def slide_results_hit_rate(prs, results):
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, H, C_LIGHT)
    header_bar(s, 'Results: Cache Hit Rate', 'Massive improvement from 0% to 99.54%')

    # Big stat cards
    cards = [
        ('0%',     'Baseline',    C_RED,   Inches(0.4)),
        ('99.54%', 'Cache Only',  C_BLUE,  Inches(4.8)),
        ('99.54%', 'AI Pre-fetch',C_GREEN, Inches(9.2)),
    ]
    for val, label, color, x in cards:
        add_rect(s, x, Inches(1.5), Inches(3.8), Inches(2.5), color)
        add_text(s, val, x+Inches(0.1), Inches(1.6), Inches(3.6), Inches(1.4),
                 font_size=48, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(s, label, x+Inches(0.1), Inches(3.0), Inches(3.6), Inches(0.5),
                 font_size=16, color=C_WHITE, align=PP_ALIGN.CENTER)

    chart_path = os.path.join(CHART_DIR, 'cache_hit_rate.png')
    add_image(s, chart_path, Inches(2.5), Inches(4.2), Inches(8.0))

    add_text(s, 'Cache layer alone eliminates almost all remote DB calls',
             Inches(0.4), Inches(6.7), Inches(12.5), Inches(0.4),
             font_size=14, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)


def slide_results_response_time(prs, results):
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, H, C_LIGHT)
    header_bar(s, 'Results: Query Response Time', 'Up to 126× speedup with cache')

    speedup_cache = results['baseline']['avg_response_ms'] / results['cache_only']['avg_response_ms']
    speedup_ai    = results['baseline']['avg_response_ms'] / results['ai_prefetch']['avg_response_ms']

    metrics = [
        (f"{results['baseline']['avg_response_ms']} ms",    'Baseline Avg',     C_RED,   Inches(0.4)),
        (f"{results['cache_only']['avg_response_ms']} ms",  'Cache Only Avg',   C_BLUE,  Inches(3.7)),
        (f"{results['ai_prefetch']['avg_response_ms']} ms", 'AI Pre-fetch Avg', C_GREEN, Inches(7.0)),
        (f'{speedup_cache:.0f}×',                           'Cache Speedup',    C_DARK,  Inches(10.3)),
    ]
    for val, label, color, x in metrics:
        add_rect(s, x, Inches(1.5), Inches(2.8), Inches(1.6), color)
        add_text(s, val, x+Inches(0.1), Inches(1.6), Inches(2.6), Inches(0.9),
                 font_size=28, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(s, label, x+Inches(0.1), Inches(2.5), Inches(2.6), Inches(0.45),
                 font_size=12, color=C_WHITE, align=PP_ALIGN.CENTER)

    chart_path = os.path.join(CHART_DIR, 'response_time_comparison.png')
    add_image(s, chart_path, Inches(0.5), Inches(3.4), Inches(12.3))


def slide_analysis(prs):
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, H, C_LIGHT)
    header_bar(s, 'Analysis: When AI Pre-fetch Makes the Difference',
               'Cold cache scenario — before the cache has warmed up')

    chart_path = os.path.join(CHART_DIR, 'nb_cold_cache_simulation.png')
    add_image(s, chart_path, Inches(0.3), Inches(1.3), Inches(8.5))

    add_rect(s, Inches(9.0), Inches(1.3), Inches(3.9), Inches(5.7), C_DARK)
    add_text(s, 'Key Insight', Inches(9.1), Inches(1.4), Inches(3.7), Inches(0.5),
             font_size=17, bold=True, color=C_BLUE)
    insights = [
        'In benchmark: both cache strategies achieve ~99.5% hit rate due to high query repetition',
        'In real cold-cache: AI Pre-fetch warms cache BEFORE queries arrive',
        'Benefit = eliminating the first N cache misses',
        'Remote queries (Site B) go from 150ms to <1ms',
        'AI adds ~0.4ms overhead per query (background thread)',
    ]
    bullet_box(s, insights, Inches(9.1), Inches(1.95), Inches(3.7), Inches(4.8),
               font_size=13, color=C_WHITE)


def slide_conclusion(prs, results):
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, H, C_LIGHT)
    header_bar(s, 'Conclusions', 'Query Predictor — Self-Tuning Distributed DB')

    rows = [
        ['Metric',                  'Baseline',   'Cache Only',       'AI Pre-fetch'],
        ['Cache Hit Rate',          '0%',         '99.54%',           '99.54%'],
        ['Avg Response Time',       '90.1 ms',    '0.71 ms',          '1.07 ms'],
        ['Speedup vs Baseline',     '1×',         '126×',             '84×'],
        ['Cold Cache Advantage',    '—',          'No',               'Yes ✓'],
        ['Prediction Required',     'No',         'No',               'Yes (LSTM)'],
    ]
    # Draw table
    col_widths = [Inches(2.8), Inches(2.0), Inches(2.4), Inches(2.4)]
    col_colors = [C_DARK, C_RED, C_BLUE, C_GREEN]
    x_start = Inches(0.4)
    y_start = Inches(1.5)

    for c, (header, color, w) in enumerate(zip(rows[0], col_colors, col_widths)):
        xp = x_start + sum(col_widths[:c])
        add_rect(s, xp, y_start, w, Inches(0.45), color)
        add_text(s, header, xp+Inches(0.05), y_start+Inches(0.02), w-Inches(0.1), Inches(0.4),
                 font_size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    for r, row in enumerate(rows[1:]):
        bg = RGBColor(0xE8,0xF4,0xFF) if r % 2 == 0 else C_WHITE
        row_y = y_start + Inches(0.45) + r * Inches(0.52)
        for c, (cell, w) in enumerate(zip(row, col_widths)):
            xp = x_start + sum(col_widths[:c])
            add_rect(s, xp, row_y, w, Inches(0.52), bg)
            add_text(s, cell, xp+Inches(0.05), row_y+Inches(0.05), w-Inches(0.1), Inches(0.42),
                     font_size=13, color=C_DARK, align=PP_ALIGN.CENTER)

    takeaways = [
        'Cache layer is the biggest single improvement in distributed DB performance',
        'ML pre-fetching converts reactive cache into a proactive system',
        'Markov Chain — fast & simple; LSTM — captures complex session dependencies',
        'System architecture scales to real distributed deployments (Docker → Kubernetes)',
    ]
    bullet_box(s, takeaways, Inches(0.4), Inches(5.35), Inches(12.5), Inches(1.8),
               font_size=14, color=C_DARK, bullet='✓')


def slide_qa(prs):
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, H, C_DARK)
    add_rect(s, 0, Inches(3.2), W, Inches(0.06), C_BLUE)

    add_text(s, 'Thank You', Inches(0.5), Inches(0.8), Inches(12), Inches(1.2),
             font_size=60, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(s, 'Q & A', Inches(0.5), Inches(2.0), Inches(12), Inches(0.8),
             font_size=36, color=C_BLUE, align=PP_ALIGN.CENTER)

    details = [
        'Source Code    :  database-distribute/',
        'Dataset        :  data/query_history_logs.csv  (500K rows)',
        'Models         :  models/markov_model.pkl  +  models/lstm_model.pt',
        'Charts         :  results/charts/',
        'Run pipeline   :  python main.py all',
    ]
    bullet_box(s, details, Inches(2.5), Inches(4.0), Inches(8.5), Inches(2.8),
               font_size=15, color=RGBColor(0xAA,0xCC,0xEE), bullet='▸')


# ── Main ──────────────────────────────────────────────────────────────────────
def main():
    with open(RESULTS_PATH) as f:
        results = json.load(f)

    prs = new_prs()

    print('Building slides...')
    slide_title(prs)             ; print('  [1/9] Title')
    slide_problem(prs)           ; print('  [2/9] Problem Statement')
    slide_architecture(prs)      ; print('  [3/9] Architecture')
    slide_dataset(prs)           ; print('  [4/9] Dataset')
    slide_ml_models(prs)         ; print('  [5/9] ML Models')
    slide_benchmark_setup(prs)   ; print('  [6/9] Benchmark Setup')
    slide_results_hit_rate(prs, results)      ; print('  [7/9] Results: Hit Rate')
    slide_results_response_time(prs, results) ; print('  [8/9] Results: Response Time')
    slide_analysis(prs)          ; print('  [9/9] Analysis')
    slide_conclusion(prs, results); print(' [10/10] Conclusion')
    slide_qa(prs)                ; print(' [11/11] Q&A')

    os.makedirs('results', exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(f'\nSaved -> {OUTPUT_PATH}')


if __name__ == '__main__':
    main()
